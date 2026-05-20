from pathlib import Path
import shutil
import subprocess
from zipfile import ZipFile, ZIP_DEFLATED

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PptInches, Pt as PptPt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image as RLImage,
    KeepTogether,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "entrega-final"
EVID = OUT / "evidencias"
QA = OUT / "qa"

TEAM = [
    "Juan Camilo Ballesteros Sierra",
    "Luis Felipe Murillo Matallana",
    "Juan Sebastián Delgado",
    "Daniela Castro Quiñones",
]

FILES = {
    "dashboard": EVID / "01_zabbix_dashboard.png",
    "hosts": EVID / "02_zabbix_hosts.png",
    "latest": EVID / "03_zabbix_latest_data.png",
    "problems": EVID / "04_zabbix_problems.png",
    "mailhog": EVID / "05_mailhog_alertas.png",
    "failure": EVID / "06_zabbix_falla_web_activa.png",
    "mailhog_failure": EVID / "07_mailhog_falla_web.png",
}


def ensure_dirs():
    OUT.mkdir(exist_ok=True)
    EVID.mkdir(exist_ok=True)
    QA.mkdir(exist_ok=True)


def set_cell_shading(cell, fill):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_cell_text(cell, text, bold=False):
    cell.text = ""
    p = cell.paragraphs[0]
    run = p.add_run(text)
    run.bold = bold
    run.font.size = Pt(8.5)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def set_two_columns(section):
    sect_pr = section._sectPr
    cols = sect_pr.xpath("./w:cols")
    if cols:
        cols = cols[0]
    else:
        cols = OxmlElement("w:cols")
        sect_pr.append(cols)
    cols.set(qn("w:num"), "2")
    cols.set(qn("w:space"), "360")


def style_doc(doc):
    for section in doc.sections:
        section.top_margin = Inches(0.7)
        section.bottom_margin = Inches(0.7)
        section.left_margin = Inches(0.62)
        section.right_margin = Inches(0.62)
    normal = doc.styles["Normal"]
    normal.font.name = "Times New Roman"
    normal.font.size = Pt(9)
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
    for name in ["Heading 1", "Heading 2"]:
        style = doc.styles[name]
        style.font.name = "Times New Roman"
        style.font.bold = True
        style.font.color.rgb = RGBColor(0x11, 0x32, 0x3d)
        style.font.size = Pt(10 if name == "Heading 1" else 9)


def add_heading(doc, text):
    p = doc.add_paragraph()
    p.style = doc.styles["Heading 1"]
    p.paragraph_format.space_before = Pt(5)
    p.paragraph_format.space_after = Pt(2)
    p.add_run(text)


def add_body(doc, text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    p.paragraph_format.first_line_indent = Inches(0.18)
    p.paragraph_format.space_after = Pt(2)
    p.add_run(text)


def add_figure(doc, path, caption, width=3.05):
    if not path.exists():
        return
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(path), width=Inches(width))
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cap.add_run(caption)
    run.italic = True
    run.font.size = Pt(8)


def create_docx():
    doc = Document()
    style_doc(doc)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("Proyecto 7: Monitoreo de infraestructura con Zabbix")
    run.bold = True
    run.font.name = "Times New Roman"
    run.font.size = Pt(16)

    authors = doc.add_paragraph()
    authors.alignment = WD_ALIGN_PARAGRAPH.CENTER
    authors.add_run("\n".join(TEAM)).font.size = Pt(10)

    abstract = doc.add_paragraph()
    abstract.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    r = abstract.add_run("Resumen - ")
    r.bold = True
    abstract.add_run(
        "Este proyecto implementa una plataforma de monitoreo de infraestructura con Zabbix 6.x desplegada completamente en contenedores Docker. "
        "La solución monitorea servicios web, base de datos, DNS y FTP, registra métricas de disponibilidad y recursos, define triggers de falla y valida el envío de alertas mediante MailHog. "
        "Adicionalmente se publicó el entorno en una VPS con HTTPS y se implementó una aplicación real con frontend, backend Node.js, MariaDB, gráficas operativas, SLO, exporter de métricas, matriz de cumplimiento y pruebas de carga con Artillery."
    )

    keywords = doc.add_paragraph()
    keywords.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    r = keywords.add_run("Palabras clave - ")
    r.bold = True
    keywords.add_run("Zabbix, Docker Compose, monitoreo, alertas, MailHog, Artillery, infraestructura.")

    doc.add_section(WD_SECTION.CONTINUOUS)
    set_two_columns(doc.sections[-1])

    add_heading(doc, "I. Introducción")
    add_body(
        doc,
        "Las infraestructuras telemáticas requieren observabilidad para detectar fallas, analizar disponibilidad y reducir tiempos de indisponibilidad. "
        "Zabbix centraliza métricas, eventos y alertas, por lo que es adecuado para demostrar monitoreo de servicios en un entorno reproducible con Docker. "
        "El alcance se amplió con un portal público que permite observar carga, telemetría y cumplimiento del enunciado durante la sustentación.",
    )
    add_body(
        doc,
        "El objetivo fue construir una solución que no solo evidenciara puertos abiertos, sino que permitiera comprobar estado de hosts, recursos, servicios, alertas, históricos y comportamiento bajo carga controlada.",
    )

    add_heading(doc, "II. Contexto del problema")
    add_body(
        doc,
        "Una red compuesta por servicios HTTP, base de datos, DNS y FTP puede presentar fallas por caída de procesos, saturación de recursos o pérdida de conectividad. "
        "Sin una plataforma de monitoreo, la detección depende de revisión manual. El proyecto busca evidenciar estado en tiempo real, historial, alertas automáticas y respuesta del sistema bajo carga controlada.",
    )
    add_body(
        doc,
        "El escenario simulado representa una infraestructura común: un usuario consume un portal web, el portal depende de base de datos, la red requiere resolución DNS y existen servicios auxiliares de transferencia. Una falla en cualquiera de estos puntos debe quedar visible para operación.",
    )

    add_heading(doc, "III. Alternativas de solución")
    add_body(
        doc,
        "Se evaluaron alternativas como Nagios, Prometheus, Datadog y Zabbix. Zabbix fue seleccionado porque integra servidor, agentes, frontend web, triggers, dashboards y notificaciones sin depender de una plataforma externa paga.",
    )
    alt_table = doc.add_table(rows=1, cols=3)
    alt_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for cell, text in zip(alt_table.rows[0].cells, ["Herramienta", "Fortaleza", "Limitación para el proyecto"]):
        set_cell_text(cell, text, True)
        set_cell_shading(cell, "D9EAD3")
    for row in [
        ("Nagios", "Disponibilidad y plugins", "Menos integrado para dashboards e históricos modernos"),
        ("Prometheus", "Métricas y series de tiempo", "Requiere integrar Alertmanager/Grafana"),
        ("Datadog", "Suite gestionada completa", "Servicio externo pago"),
        ("Zabbix", "Agentes, triggers, UI, alertas y API", "Requiere configuración inicial cuidadosa"),
    ]:
        cells = alt_table.add_row().cells
        for cell, text in zip(cells, row):
            set_cell_text(cell, text)

    add_heading(doc, "IV. Diseño de la solución")
    add_body(
        doc,
        "La arquitectura incluye Zabbix Server, PostgreSQL, Zabbix Web, MailHog y cuatro servicios monitoreados. Todos los componentes se conectan mediante la red Docker proyecto7-monitoring. "
        "Los agentes reportan disponibilidad y Zabbix Server ejecuta checks simples para validar el estado de los puertos de servicio. En la VPS, Caddy publica Zabbix, MailHog y el portal web mediante subdominios HTTPS.",
    )
    add_body(
        doc,
        "El servicio web monitoreado se diseñó como una aplicación observable: expone salud, resumen operativo, estado de base de datos, telemetría, incidentes, carga controlada, gráficas, SLO y matriz de cumplimiento. Esto permite observar comportamiento funcional, no solo conectividad.",
    )

    add_figure(doc, FILES["dashboard"], "Fig. 1. Dashboard principal del proyecto en Zabbix.")

    table = doc.add_table(rows=1, cols=3)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr = table.rows[0].cells
    for cell, text in zip(hdr, ["Host", "Servicio", "Check"]):
        set_cell_text(cell, text, True)
        set_cell_shading(cell, "D9EAD3")
    rows = [
        ("web-host", "Node.js web", "HTTP puerto 80"),
        ("db-host", "MariaDB", "TCP puerto 3306"),
        ("dns-host", "CoreDNS", "TCP puerto 53"),
        ("ftp-host", "VSFTPD", "FTP puerto 21"),
    ]
    for row in rows:
        cells = table.add_row().cells
        for cell, text in zip(cells, row):
            set_cell_text(cell, text)

    add_figure(doc, FILES["hosts"], "Fig. 2. Hosts monitoreados registrados en Zabbix.")

    add_heading(doc, "V. Implementación")
    add_body(
        doc,
        "El despliegue se ejecuta con Docker Compose. El archivo .env centraliza credenciales, puertos y zona horaria. "
        "El servidor Zabbix usa una imagen personalizada construida desde docker/zabbix-server/Dockerfile. "
        "También se montan archivos de configuración Zabbix como volumen para el servidor y los agentes. "
        "El script de aprovisionamiento usa la API JSON-RPC de Zabbix para crear el grupo de hosts, items, triggers, dashboard, escenario web y configuración de correo hacia MailHog. "
        "El servicio web expone endpoints JSON, /metrics, /api/charts y /api/compliance para demostrar monitoreo avanzado.",
    )
    add_body(
        doc,
        "En la VPS se usa docker-compose.vps.yml para mantener internos los puertos de servicio y publicar solamente las interfaces necesarias mediante Caddy. MailHog se protege con un gate de autenticación y Zabbix conserva los históricos en PostgreSQL.",
    )
    add_figure(doc, FILES["latest"], "Fig. 3. Datos recientes de disponibilidad y métricas.")

    add_heading(doc, "VI. Pruebas")
    add_body(
        doc,
        "Se validaron cuatro escenarios mínimos: dashboard en tiempo real, simulación de caída del servicio web, envío de alertas a MailHog y consulta de datos históricos. "
        "También se ejecutaron pruebas con Artillery contra frontend, API, base de datos y endpoints de carga; la auditoría automática revisó Compose, endpoints públicos, matriz de cumplimiento y objetos principales de Zabbix.",
    )
    add_body(
        doc,
        "La prueba smoke en producción generó 96 solicitudes, 0 usuarios fallidos y p95 aproximado de 23.8 ms. La auditoría automática reportó 27 validaciones OK y 0 fallidas, incluyendo endpoints públicos, Zabbix API y matriz de cumplimiento.",
    )
    add_figure(doc, FILES["failure"], "Fig. 4. Problema activo durante la simulación de caída.")
    add_figure(doc, FILES["mailhog_failure"], "Fig. 5. Correos de alerta y recuperación recibidos en MailHog.")

    add_heading(doc, "VII. Discusión")
    add_body(
        doc,
        "Los checks de servicio regresan 1 cuando el puerto responde y 0 cuando no responde. Los triggers permiten convertir cambios de estado en eventos visibles y notificaciones. "
        "MailHog facilita probar el flujo de correo sin exponer cuentas reales y el canal SMTP real del dominio demuestra escalamiento externo. "
        "Las pruebas de carga relacionan tráfico, latencia, SLO, escrituras en MariaDB y gráficas históricas.",
    )
    add_body(
        doc,
        "La principal diferencia frente a una instalación básica es que la aplicación monitoreada produce datos propios. Esto permite discutir degradación, escritura en base de datos, rutas más consultadas y disponibilidad calculada, aspectos más cercanos a una operación real.",
    )

    add_heading(doc, "VIII. Conclusiones")
    add_body(
        doc,
        "La solución cumple los requerimientos del Proyecto 7: infraestructura en contenedores, mínimo cuatro hosts monitoreados, Zabbix Server con base de datos, frontend web, triggers, dashboards y alertas. "
        "El despliegue público, backend transaccional, Artillery, exporter /metrics y matriz /api/compliance agregan evidencia para defender el proyecto por encima de los requisitos mínimos.",
    )

    add_heading(doc, "Referencias")
    refs = [
        "Zabbix Documentation, https://www.zabbix.com/documentation",
        "Docker Documentation, https://docs.docker.com",
        "MailHog, https://github.com/mailhog/MailHog",
        "CoreDNS, https://coredns.io",
        "Artillery Documentation, https://www.artillery.io/docs",
        "Caddy Documentation, https://caddyserver.com/docs",
    ]
    for ref in refs:
        p = doc.add_paragraph(ref)
        p.paragraph_format.left_indent = Inches(0.12)
        p.paragraph_format.space_after = Pt(1)

    out = OUT / "Informe_IEEE_Proyecto7_Zabbix.docx"
    doc.save(out)
    return out


def rl_img(path, width=6.2 * inch):
    img = Image.open(path)
    ratio = img.height / img.width
    return RLImage(str(path), width=width, height=width * ratio)


def create_pdf():
    out = OUT / "Informe_IEEE_Proyecto7_Zabbix.pdf"
    doc = SimpleDocTemplate(
        str(out),
        pagesize=letter,
        rightMargin=0.62 * inch,
        leftMargin=0.62 * inch,
        topMargin=0.7 * inch,
        bottomMargin=0.7 * inch,
    )
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle("TitleCenter", parent=styles["Title"], alignment=TA_CENTER, fontName="Times-Bold", fontSize=16, leading=18))
    styles.add(ParagraphStyle("Author", parent=styles["Normal"], alignment=TA_CENTER, fontName="Times-Roman", fontSize=9, leading=11))
    styles.add(ParagraphStyle("IEEEBody", parent=styles["Normal"], alignment=TA_JUSTIFY, fontName="Times-Roman", fontSize=9, leading=11, firstLineIndent=12))
    styles.add(ParagraphStyle("IEEEHead", parent=styles["Heading2"], fontName="Times-Bold", fontSize=10, leading=12, textColor=colors.HexColor("#11323d")))
    styles.add(ParagraphStyle("Caption", parent=styles["Normal"], alignment=TA_CENTER, fontName="Times-Italic", fontSize=8, leading=9))

    story = [
        Paragraph("Proyecto 7: Monitoreo de infraestructura con Zabbix", styles["TitleCenter"]),
        Paragraph("<br/>".join(TEAM), styles["Author"]),
        Spacer(1, 0.12 * inch),
        Paragraph("<b>Resumen - </b>Este proyecto implementa una plataforma de monitoreo con Zabbix 6.x, Docker Compose, cuatro servicios monitoreados y alertas validadas con MailHog. La solución se publicó en VPS con HTTPS e incorpora backend Node.js, MariaDB, gráficas, SLO, exporter de métricas, matriz de cumplimiento y pruebas Artillery.", styles["IEEEBody"]),
        Paragraph("<b>Palabras clave - </b>Zabbix, Docker, monitoreo, alertas, MailHog, Artillery.", styles["IEEEBody"]),
        Spacer(1, 0.12 * inch),
    ]
    sections = [
        ("I. Introducción", "Las infraestructuras telemáticas requieren observabilidad para detectar fallas, analizar disponibilidad y reducir tiempos de indisponibilidad. El objetivo fue construir una solución que evidenciara estado de hosts, recursos, servicios, alertas, históricos y comportamiento bajo carga controlada."),
        ("II. Contexto del problema", "La red evaluada contiene servicios HTTP, base de datos, DNS y FTP. Sin monitoreo, la detección depende de revisión manual y no queda historial para análisis. La solución debe responder si el host está activo, si el servicio responde, si hay consumo anormal y si la alerta llega por correo."),
        ("III. Alternativas de solución", "Se compararon Nagios, Prometheus, Datadog y Zabbix. Zabbix fue elegido porque integra agentes, servidor, frontend, base de datos, triggers, dashboards, web scenarios, media types y API de aprovisionamiento en una sola plataforma abierta."),
        ("IV. Diseño de la solución", "La arquitectura usa Zabbix Server, PostgreSQL, Zabbix Web, MailHog, cuatro servicios monitoreados y agentes Zabbix. El portal público agrega backend Node.js, MariaDB, endpoints JSON, /metrics, /api/charts, /api/live y /api/compliance. En la VPS, Caddy publica subdominios HTTPS."),
    ]
    for title, body in sections:
        story += [Paragraph(title, styles["IEEEHead"]), Paragraph(body, styles["IEEEBody"])]

    alt = [["Herramienta", "Fortaleza", "Limitación"], ["Nagios", "Checks y plugins", "Menos integrado para gráficas modernas"], ["Prometheus", "Métricas y series de tiempo", "Requiere Alertmanager/Grafana"], ["Datadog", "Suite gestionada", "Servicio externo pago"], ["Zabbix", "Agentes, UI, triggers y alertas", "Mayor configuración inicial"]]
    alt_tbl = Table(alt, colWidths=[1.2 * inch, 2.0 * inch, 2.5 * inch])
    alt_tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#D9EAD3")),
        ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#9CA3AF")),
        ("FONTNAME", (0, 0), (-1, 0), "Times-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 7.5),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]))
    story += [Spacer(1, 0.06 * inch), alt_tbl, Spacer(1, 0.08 * inch)]

    data = [["Host", "Servicio", "Check"], ["web-host", "Node.js web", "HTTP 80"], ["db-host", "MariaDB", "TCP 3306"], ["dns-host", "CoreDNS", "TCP 53"], ["ftp-host", "VSFTPD", "FTP 21"]]
    tbl = Table(data, colWidths=[1.4 * inch, 1.4 * inch, 2.0 * inch])
    tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#D9EAD3")),
        ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#9CA3AF")),
        ("FONTNAME", (0, 0), (-1, 0), "Times-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]))
    story += [Spacer(1, 0.08 * inch), tbl, Spacer(1, 0.12 * inch)]
    for key, caption in [
        ("dashboard", "Fig. 1. Dashboard principal del proyecto en Zabbix."),
        ("hosts", "Fig. 2. Hosts monitoreados en Zabbix."),
        ("latest", "Fig. 3. Latest data con checks de disponibilidad."),
        ("failure", "Fig. 4. Problema activo durante la caída del servicio web."),
        ("mailhog_failure", "Fig. 5. Alertas recibidas en MailHog."),
    ]:
        if FILES[key].exists():
            story += [KeepTogether([rl_img(FILES[key], 5.75 * inch), Paragraph(caption, styles["Caption"])])]

    story += [
        Paragraph("V. Implementación", styles["IEEEHead"]),
        Paragraph("El despliegue se empaqueta con docker-compose.yml para entorno local y docker-compose.vps.yml para publicación. Zabbix Server usa una imagen personalizada y monta configuraciones como volumen. El script provision_zabbix.py crea grupo, hosts, items, triggers, dashboard, media type y web scenario mediante API JSON-RPC.", styles["IEEEBody"]),
        Paragraph("La aplicación web expone endpoints para salud, resumen, base de datos, telemetría, incidentes, carga controlada, gráficas, SLO, cumplimiento y métricas estilo Prometheus. Esto permite observar comportamiento funcional, no solo puertos abiertos.", styles["IEEEBody"]),
        Paragraph("La configuración final deja trazabilidad entre requisito, archivo y evidencia: Compose define la topología, los Dockerfile y volúmenes prueban personalización, los scripts automatizan Zabbix y las capturas muestran dashboard, hosts, latest data, problemas y correos.", styles["IEEEBody"]),
        Paragraph("VI. Pruebas", styles["IEEEHead"]),
        Paragraph("Se validaron dashboard en tiempo real, caída controlada, envío de alertas, métricas históricas y carga con Artillery. En producción, el smoke test generó 96 solicitudes, 0 usuarios fallidos y p95 aproximado de 23.8 ms. La auditoría automática reportó 27 validaciones OK y 0 fallidas.", styles["IEEEBody"]),
        Paragraph("Durante la sustentación se puede ejecutar una prueba en vivo con artillery-live-demo.yml y detener temporalmente web-service. Con esto el profesor observa tráfico real en el portal, detección de falla en Zabbix, alerta en MailHog y recuperación al reiniciar el servicio.", styles["IEEEBody"]),
        Paragraph("VII. Discusión de las pruebas", styles["IEEEHead"]),
        Paragraph("La prueba de caída confirma el flujo completo: falla, problema en Zabbix, alerta en MailHog, recuperación y registro histórico. Las pruebas de carga permiten diferenciar caída total de degradación por tráfico, mientras /metrics y /api/db/status amplían el monitoreo más allá de conectividad básica.", styles["IEEEBody"]),
        Paragraph("El uso de MailHog evita enviar correos reales durante las pruebas, pero el proyecto conserva una ruta SMTP del dominio para explicar cómo se escalaría la notificación en un ambiente de producción.", styles["IEEEBody"]),
        Paragraph("VIII. Conclusiones", styles["IEEEHead"]),
        Paragraph("La solución cumple los requerimientos: despliegue dockerizado, mínimo cuatro hosts monitoreados, dashboards, triggers, alertas y evidencia histórica. El portal público, backend transaccional, SLO, exporter /metrics, Artillery y matriz /api/compliance agregan valor para la sustentación y trazabilidad frente a la rúbrica.", styles["IEEEBody"]),
        Paragraph("Referencias", styles["IEEEHead"]),
        Paragraph("[1] Zabbix Documentation. [2] Docker Documentation. [3] MailHog GitHub. [4] Artillery Documentation. [5] Caddy Documentation.", styles["IEEEBody"]),
    ]
    doc.build(story)
    return out


def add_text(slide, text, x, y, w, h, size=24, bold=False, color=(22, 41, 52), align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = PptPt(size)
    r.font.bold = bold
    r.font.color.rgb = PptRGB(*color)
    return box


def add_bullets(slide, bullets, x, y, w, h, size=20):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = PptPt(size)
        p.font.color.rgb = PptRGB(35, 45, 55)
        p.space_after = PptPt(8)
    return box


def add_screenshot(slide, path, x, y, w, h=None):
    if not path.exists():
        return None
    with Image.open(path) as img:
        ratio = img.height / img.width
    if h is None:
        h = w * ratio
    return slide.shapes.add_picture(str(path), PptInches(x), PptInches(y), width=PptInches(w), height=PptInches(h))


def add_footer(slide, num):
    add_text(slide, f"Proyecto 7 - Zabbix | {num}", 0.45, 7.12, 4.0, 0.25, size=8, color=(93, 107, 122))


def create_pptx():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]
    total_slides = 17

    INK = (10, 24, 32)
    TEAL = (34, 197, 170)
    CYAN = (56, 189, 248)
    LIME = (163, 230, 53)
    AMBER = (245, 158, 11)
    RED = (239, 68, 68)
    PAPER = (247, 250, 247)
    MUTED = (92, 105, 117)
    LINE = (191, 204, 214)
    WHITE = (255, 255, 255)

    def rgb(color):
        return PptRGB(*color)

    def bg(slide, color=PAPER):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(color)

    def rect(slide, x, y, w, h, fill, line=None, radius=False):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shp = slide.shapes.add_shape(shape_type, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
        if line:
            shp.line.color.rgb = rgb(line)
            shp.line.width = PptPt(0.8)
        else:
            shp.line.fill.background()
        return shp

    def line(slide, x1, y1, x2, y2, color=LINE, width=1.4):
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            PptInches(x1),
            PptInches(y1),
            PptInches(x2),
            PptInches(y2),
        )
        connector.line.color.rgb = rgb(color)
        connector.line.width = PptPt(width)
        return connector

    def text(slide, value, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
        box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
        tf = box.text_frame
        tf.clear()
        tf.margin_left = PptInches(0)
        tf.margin_right = PptInches(0)
        tf.margin_top = PptInches(0)
        tf.margin_bottom = PptInches(0)
        tf.word_wrap = True
        tf.vertical_anchor = valign
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = value
        r.font.name = "Aptos Display" if size >= 26 else "Aptos"
        r.font.size = PptPt(size)
        r.font.bold = bold
        r.font.color.rgb = rgb(color)
        return box

    def kicker(slide, label, num, dark=False):
        color = TEAL if dark else INK
        text(slide, f"{num:02d}", 0.55, 0.42, 0.42, 0.22, 8, color, True)
        line(slide, 1.03, 0.52, 1.78, 0.52, color, 1.1)
        text(slide, label.upper(), 1.9, 0.39, 3.4, 0.3, 8, color, True)

    def title(slide, label, claim, num, dark=False):
        kicker(slide, label, num, dark)
        text(slide, claim, 0.55, 0.78, 8.2, 0.95, 30, WHITE if dark else INK, True)

    def footer(slide, num, dark=False):
        color = (142, 159, 170) if dark else MUTED
        line(slide, 0.55, 7.03, 12.75, 7.03, (42, 66, 77) if dark else LINE, 0.7)
        text(slide, "Proyecto 7 - Monitoreo de infraestructura con Zabbix", 0.55, 7.13, 5.2, 0.18, 7.5, color)
        text(slide, f"{num:02d}/{total_slides}", 12.0, 7.13, 0.7, 0.18, 7.5, color, align=PP_ALIGN.RIGHT)

    def screenshot(slide, path, x, y, w, h=None, border=True):
        if not path.exists():
            return None
        if h is None:
            with Image.open(path) as img:
                h = w * img.height / img.width
        pic = slide.shapes.add_picture(str(path), PptInches(x), PptInches(y), width=PptInches(w), height=PptInches(h))
        if border:
            box = rect(slide, x - 0.03, y - 0.03, w + 0.06, h + 0.06, (255, 255, 255), LINE)
            slide.shapes._spTree.remove(box._element)
            slide.shapes._spTree.insert(2, box._element)
        return pic

    def metric(slide, value, label, x, y, w=1.75, accent=TEAL, dark=False):
        rect(slide, x, y, w, 0.78, (18, 43, 54) if dark else WHITE, (48, 73, 84) if dark else LINE, True)
        text(slide, value, x + 0.16, y + 0.12, w - 0.28, 0.28, 18, accent, True)
        text(slide, label, x + 0.16, y + 0.45, w - 0.28, 0.2, 7.8, (202, 213, 219) if dark else MUTED, True)

    def label_box(slide, heading, body, x, y, w, h, accent=TEAL, dark=False):
        rect(slide, x, y, w, h, (16, 40, 51) if dark else WHITE, (48, 73, 84) if dark else LINE, True)
        rect(slide, x, y, 0.08, h, accent)
        text(slide, heading, x + 0.25, y + 0.16, w - 0.42, 0.24, 12, accent, True)
        text(slide, body, x + 0.25, y + 0.47, w - 0.42, h - 0.56, 11, (223, 232, 238) if dark else INK)

    def node(slide, label, sub, x, y, w=1.45, h=0.72, fill=WHITE, accent=TEAL):
        rect(slide, x, y, w, h, fill, LINE, True)
        text(slide, label, x + 0.12, y + 0.12, w - 0.24, 0.22, 11, INK, True, PP_ALIGN.CENTER)
        text(slide, sub, x + 0.12, y + 0.39, w - 0.24, 0.18, 7.2, MUTED, align=PP_ALIGN.CENTER)
        line(slide, x + 0.15, y + h - 0.1, x + w - 0.15, y + h - 0.1, accent, 2.2)

    # 1 cover
    s = prs.slides.add_slide(blank)
    bg(s, INK)
    rect(s, 8.25, 0.0, 5.08, 7.5, (15, 40, 52))
    screenshot(s, FILES["dashboard"], 7.35, 0.76, 5.45, 3.55, True)
    text(s, "PROYECTO 7", 0.7, 0.72, 2.3, 0.25, 9, TEAL, True)
    text(s, "Monitoreo de infraestructura\ncon Zabbix", 0.68, 1.23, 6.2, 1.35, 38, WHITE, True)
    text(s, "Zabbix 6.x monitoreando servicios Docker, alertas, backend real, métricas, SLO y pruebas de carga en vivo.", 0.72, 3.02, 5.8, 0.65, 16, (205, 220, 228))
    metric(s, "4", "hosts monitoreados", 0.72, 4.15, 1.55, TEAL, True)
    metric(s, "7", "páginas IEEE", 2.45, 4.15, 1.55, CYAN, True)
    metric(s, "96", "requests smoke", 4.18, 4.15, 1.55, LIME, True)
    metric(s, "0", "fallas auditoría", 5.91, 4.15, 1.55, AMBER, True)
    text(s, "\n".join(TEAM), 0.72, 5.72, 6.25, 0.9, 12, (221, 232, 238))
    footer(s, 1, True)

    # 2 problem
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "problema", "El riesgo real es enterarse tarde de una caída.", 2)
    text(s, "La infraestructura evaluada mezcla web, base de datos, DNS y FTP. Si un componente cae, el usuario percibe degradación antes de que el equipo tenga una causa clara.", 0.62, 1.85, 5.4, 0.72, 14, MUTED)
    stages = [("Falla", "servicio cae", RED), ("Usuario", "reporta tarde", AMBER), ("Operación", "revisa manual", CYAN), ("Histórico", "no existe", TEAL)]
    for i, (name, sub, accent) in enumerate(stages):
        x = 0.85 + i * 3.05
        node(s, name, sub, x, 3.55, 1.72, 0.82, WHITE, accent)
        if i < 3:
            line(s, x + 1.72, 3.96, x + 2.68, 3.96, LINE, 2)
    label_box(s, "Riesgo operativo", "No basta saber que el servidor está encendido; se debe medir servicio, recursos, alertas e histórico.", 7.2, 1.75, 4.95, 1.35, TEAL)
    label_box(s, "Criterio de evaluación", "La rúbrica pide diseño, proceso de validación y construcción funcional. Por eso la demo debe probar el flujo completo.", 7.2, 4.55, 4.95, 1.35, AMBER)
    footer(s, 2)

    # 3 objective
    s = prs.slides.add_slide(blank)
    bg(s, (239, 246, 242))
    title(s, "objetivo", "Montamos algo que se puede medir, presionar y recuperar.", 3)
    label_box(s, "Mínimo del enunciado", "Zabbix Server, base de datos, frontend, MailHog, cuatro hosts, agentes, templates, triggers y dashboards.", 0.75, 2.0, 3.65, 2.0, TEAL)
    label_box(s, "Valor agregado", "Portal HTTPS con backend Node.js, MariaDB, gráficas, SLO, exporter /metrics, Artillery y auditoría reproducible.", 4.85, 2.0, 3.65, 2.0, CYAN)
    label_box(s, "Evidencia defendible", "Repositorio, README, informe IEEE de 7 páginas, PPTX, capturas, scripts y matriz /api/compliance.", 8.95, 2.0, 3.65, 2.0, LIME)
    text(s, "Tesis de la sustentación", 0.75, 5.2, 2.7, 0.25, 11, MUTED, True)
    text(s, "La gracia del proyecto está en mostrar operación: salud, carga, incidentes, alertas y recuperación.", 0.75, 5.58, 10.2, 0.48, 22, INK, True)
    footer(s, 3)

    # 4 architecture
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "arquitectura", "Zabbix observa servicios Docker.", 4)
    node(s, "Caddy", "HTTPS público", 0.85, 2.0, 1.55, 0.78, (235, 253, 246), TEAL)
    node(s, "Zabbix Web", "frontend", 3.0, 1.55, 1.65, 0.78, WHITE, CYAN)
    node(s, "Zabbix Server", "triggers + API", 3.0, 2.95, 1.65, 0.78, WHITE, TEAL)
    node(s, "PostgreSQL", "históricos", 5.25, 2.95, 1.55, 0.78, WHITE, LIME)
    node(s, "MailHog", "SMTP lab", 5.25, 1.55, 1.55, 0.78, WHITE, AMBER)
    for i, (label, sub) in enumerate([("web", "HTTP"), ("DB", "3306"), ("DNS", "53"), ("FTP", "21")]):
        node(s, label, sub, 8.0 + (i % 2) * 2.15, 1.55 + (i // 2) * 1.4, 1.55, 0.78, WHITE, [TEAL, CYAN, LIME, AMBER][i])
    line(s, 2.4, 2.39, 3.0, 1.94, LINE, 1.5)
    line(s, 2.4, 2.39, 3.0, 3.34, LINE, 1.5)
    line(s, 4.65, 3.34, 5.25, 3.34, LINE, 1.5)
    line(s, 4.65, 3.34, 5.25, 1.94, LINE, 1.5)
    line(s, 6.8, 3.34, 8.0, 1.94, LINE, 1.5)
    line(s, 6.8, 3.34, 8.0, 3.34, LINE, 1.5)
    text(s, "Infraestructura monitoreada", 8.0, 4.28, 3.7, 0.24, 10, MUTED, True, PP_ALIGN.CENTER)
    metric(s, "3", "subdominios HTTPS", 0.82, 5.35, 2.0, TEAL)
    metric(s, "12", "contenedores base", 3.15, 5.35, 2.0, CYAN)
    metric(s, "API", "aprovisionamiento", 5.48, 5.35, 2.0, LIME)
    footer(s, 4)

    # 5 inventory
    s = prs.slides.add_slide(blank)
    bg(s, (246, 248, 250))
    title(s, "inventario", "El inventario cubre web, base de datos, DNS y FTP.", 5)
    headers = ["Host", "Servicio", "Check", "Agente"]
    xs = [0.82, 3.3, 5.78, 9.05]
    widths = [2.1, 2.1, 2.85, 2.25]
    for x, w, h in zip(xs, widths, headers):
        rect(s, x, 2.0, w, 0.46, INK)
        text(s, h, x + 0.12, 2.13, w - 0.24, 0.14, 8, WHITE, True)
    rows = [
        ("web-host", "Portal Node.js", "HTTP puerto 80", "web-agent", TEAL),
        ("db-host", "MariaDB", "TCP puerto 3306", "db-agent", CYAN),
        ("dns-host", "CoreDNS", "TCP puerto 53", "dns-agent", LIME),
        ("ftp-host", "VSFTPD", "FTP puerto 21", "ftp-agent", AMBER),
    ]
    for r, row in enumerate(rows):
        y = 2.62 + r * 0.72
        for c, value in enumerate(row[:4]):
            rect(s, xs[c], y, widths[c], 0.56, WHITE, LINE)
            text(s, value, xs[c] + 0.12, y + 0.18, widths[c] - 0.24, 0.16, 9.5, INK if c != 0 else row[4], c == 0)
    screenshot(s, FILES["hosts"], 8.6, 5.33, 3.35, 1.1, True)
    text(s, "El inventario se puede defender desde Zabbix y desde Compose: cada host tiene servicio, puerto, agente y trigger asociado.", 0.85, 5.45, 6.65, 0.52, 17, INK, True)
    footer(s, 5)

    # 6 implementation
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "implementación", "Todo se puede levantar de nuevo desde el repo.", 6)
    files = [
        ("docker-compose.yml", "topología local"),
        ("docker-compose.vps.yml", "publicación VPS"),
        ("Dockerfile Zabbix", "imagen personalizada"),
        ("provision_zabbix.py", "hosts, items y triggers"),
        ("audit-project.sh", "validación reproducible"),
    ]
    for i, (name, sub) in enumerate(files):
        y = 1.95 + i * 0.72
        rect(s, 0.82, y, 4.35, 0.52, WHITE, LINE)
        text(s, name, 1.02, y + 0.13, 2.6, 0.18, 10.5, INK, True)
        text(s, sub, 3.55, y + 0.14, 1.35, 0.16, 8, MUTED, align=PP_ALIGN.RIGHT)
    for i, (value, label, accent) in enumerate([("Compose", "orquesta", TEAL), ("API", "aprovisiona", CYAN), ("Zabbix", "observa", LIME), ("MailHog", "notifica", AMBER)]):
        x = 6.0 + i * 1.62
        node(s, value, label, x, 3.0, 1.2, 0.8, WHITE, accent)
        if i < 3:
            line(s, x + 1.2, 3.4, x + 1.58, 3.4, LINE, 1.5)
    label_box(s, "Punto fuerte", "Cada cosa que se muestra sale de un archivo, un contenedor, un item de Zabbix o una evidencia.", 6.1, 5.1, 5.5, 1.1, TEAL)
    footer(s, 6)

    # 7 dashboard
    s = prs.slides.add_slide(blank)
    bg(s, INK)
    title(s, "dashboard", "El dashboard muestra el pulso del sistema.", 7, True)
    screenshot(s, FILES["dashboard"], 0.72, 1.78, 8.35, 4.72, True)
    metric(s, "CPU", "recursos", 9.55, 2.05, 1.85, TEAL, True)
    metric(s, "RAM", "memoria", 9.55, 3.05, 1.85, CYAN, True)
    metric(s, "DISCO", "capacidad", 9.55, 4.05, 1.85, LIME, True)
    metric(s, "HTTP", "servicios", 9.55, 5.05, 1.85, AMBER, True)
    footer(s, 7, True)

    # 8 latest data
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "métricas", "Latest data separa sano de caído.", 8)
    screenshot(s, FILES["latest"], 0.72, 1.75, 7.55, 4.58, True)
    label_box(s, "Checks básicos", "agent.ping, CPU, memoria y disco para observar salud del host.", 8.75, 1.95, 3.35, 1.08, TEAL)
    label_box(s, "Checks de servicio", "HTTP, MySQL/MariaDB, DNS y FTP validan disponibilidad funcional.", 8.75, 3.45, 3.35, 1.08, CYAN)
    label_box(s, "Checks avanzados", "/metrics, /api/db/status y web scenario público conectan Zabbix con la app real.", 8.75, 4.95, 3.35, 1.08, LIME)
    footer(s, 8)

    # 9 failure
    s = prs.slides.add_slide(blank)
    bg(s, (250, 247, 244))
    title(s, "caída controlada", "La caída controlada prueba el flujo completo.", 9)
    screenshot(s, FILES["failure"], 6.85, 1.7, 5.45, 3.55, True)
    steps = [("1", "stop web-service", RED), ("2", "trigger en Zabbix", AMBER), ("3", "correo MailHog", CYAN), ("4", "start + recovery", TEAL)]
    for i, (num, lbl, accent) in enumerate(steps):
        x = 0.78 + i * 1.42
        rect(s, x, 2.55, 0.72, 0.72, accent, None, True)
        text(s, num, x + 0.24, 2.74, 0.25, 0.24, 16, WHITE, True, PP_ALIGN.CENTER)
        text(s, lbl, x - 0.12, 3.46, 1.0, 0.34, 8.3, INK, True, PP_ALIGN.CENTER)
        if i < 3:
            line(s, x + 0.72, 2.91, x + 1.27, 2.91, LINE, 1.6)
    text(s, "El valor de esta prueba está en el flujo completo: no solo aparece el problema, también queda histórico y evidencia de recuperación.", 0.78, 5.55, 5.65, 0.72, 18, INK, True)
    footer(s, 9)

    # 10 mailhog
    s = prs.slides.add_slide(blank)
    bg(s, INK)
    title(s, "alertas", "MailHog captura alertas.", 10, True)
    screenshot(s, FILES["mailhog_failure"], 0.72, 1.62, 8.6, 4.86, True)
    label_box(s, "Laboratorio seguro", "Captura correos sin enviar spam a cuentas reales durante las pruebas.", 9.75, 1.85, 2.65, 1.18, TEAL, True)
    label_box(s, "Escalamiento", "Queda documentado un canal SMTP real del dominio para producción.", 9.75, 3.45, 2.65, 1.18, AMBER, True)
    footer(s, 10, True)

    # 11 artillery
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "carga", "Artillery mete presión real a la demo.", 11)
    text(s, "artillery run tests/artillery-live-demo.yml", 0.85, 2.05, 5.7, 0.32, 16, TEAL, True)
    values = [96, 0, 24]
    labels = ["requests smoke", "usuarios fallidos", "p95 aprox. ms"]
    colors_ = [TEAL, LIME, AMBER]
    for i, (val, lbl, accent) in enumerate(zip(values, labels, colors_)):
        x = 1.0 + i * 2.0
        h = 0.7 + (val / 96) * 2.15 if val else 0.22
        rect(s, x, 5.55 - h, 1.05, h, accent)
        text(s, str(val), x, 5.78, 1.05, 0.32, 18, INK, True, PP_ALIGN.CENTER)
        text(s, lbl, x - 0.2, 6.17, 1.45, 0.22, 8.3, MUTED, True, PP_ALIGN.CENTER)
    label_box(s, "Qué demuestra", "La app no solo responde: guarda telemetría, registra incidentes y expone métricas para Zabbix.", 7.15, 2.0, 4.8, 1.5, CYAN)
    label_box(s, "Cómo se ve en vivo", "Mientras corre Artillery suben requests, rutas golpeadas, cargas recientes y SLO en el portal.", 7.15, 4.1, 4.8, 1.25, TEAL)
    footer(s, 11)

    # 12 results
    s = prs.slides.add_slide(blank)
    bg(s, (242, 247, 249))
    title(s, "resultados", "El cierre se defiende con métricas y auditoría.", 12)
    metric(s, "100%", "/api/compliance", 0.85, 2.05, 2.1, TEAL)
    metric(s, "27", "validaciones OK", 3.25, 2.05, 2.1, CYAN)
    metric(s, "0", "fallas auditoría", 5.65, 2.05, 2.1, LIME)
    screenshot(s, FILES["problems"], 8.2, 1.75, 3.95, 2.7, True)
    checks = ["Dashboard en tiempo real", "Simulación de caída", "Alertas en MailHog", "Métricas históricas", "Carga Artillery", "Informe IEEE + repo"]
    for i, chk in enumerate(checks):
        y = 4.08 + (i // 2) * 0.55
        x = 0.9 + (i % 2) * 3.2
        rect(s, x, y, 0.18, 0.18, TEAL)
        text(s, chk, x + 0.3, y - 0.02, 2.6, 0.2, 9.5, INK, True)
    footer(s, 12)

    # 13 deliverables
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "entregables", "Queda repo, informe, deck y demo pública.", 13)
    items = [
        ("Informe IEEE", "entrega-final/Informe_IEEE_Proyecto7_Zabbix.pdf", TEAL),
        ("Diapositivas", "entrega-final/Presentacion_Proyecto7_Zabbix.pptx", CYAN),
        ("Repositorio", "github.com/ballesterossmartsolutionssas/Proyecto7-Zabbix", LIME),
        ("Evidencias", "entrega-final/evidencias + auditoría", AMBER),
    ]
    for i, (head, body, accent) in enumerate(items):
        x = 0.85 + (i % 2) * 5.55
        y = 2.05 + (i // 2) * 1.55
        label_box(s, head, body, x, y, 4.75, 1.05, accent)
    text(s, "Reparto: Luis Felipe abre, Juan Sebastián explica arquitectura, Daniela muestra Zabbix y Juan Camilo cierra con pruebas en vivo.", 0.9, 5.85, 10.7, 0.48, 17, INK, True)
    footer(s, 13)

    # 14 demo
    s = prs.slides.add_slide(blank)
    bg(s, INK)
    title(s, "demo en vivo", "La demo debe cerrar viendo tráfico, caída y alerta.", 14, True)
    commands = [
        ("1", "Abrir portal", "https://web-zabbix.negociocontigo.com"),
        ("2", "Generar carga", "artillery run tests/artillery-live-demo.yml"),
        ("3", "Simular caída", "docker compose -f docker-compose.vps.yml stop web-service"),
        ("4", "Validar cierre", "bash scripts/audit-project.sh"),
    ]
    for i, (num, head, cmd) in enumerate(commands):
        y = 2.05 + i * 0.95
        rect(s, 0.85, y, 0.48, 0.48, [TEAL, CYAN, AMBER, LIME][i], None, True)
        text(s, num, 1.02, y + 0.12, 0.13, 0.13, 11, INK, True, PP_ALIGN.CENTER)
        text(s, head, 1.55, y + 0.04, 2.1, 0.22, 12, WHITE, True)
        text(s, cmd, 3.55, y + 0.07, 7.4, 0.18, 10.5, (200, 216, 225))
    screenshot(s, FILES["mailhog"], 8.5, 4.45, 3.65, 1.6, True)
    text(s, "Cierre: Zabbix detecta, MailHog muestra, Artillery presiona y la auditoría confirma.", 0.9, 6.25, 7.2, 0.42, 18, TEAL, True)
    footer(s, 14, True)

    # 15 alternatives
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "alternativas", "Zabbix se eligió por equilibrio entre control, alertas y despliegue local.", 15)
    alternatives = [
        ("Nagios Core", "Checks maduros y muchos plugins.", "UI y dashboards limitados; menor automatización moderna.", AMBER),
        ("Prometheus + Grafana", "Métricas modernas y visualización fuerte.", "Requiere Alertmanager, exporters y más integración.", CYAN),
        ("Datadog", "Suite SaaS completa con APM y logs.", "Costo, dependencia externa y datos fuera del laboratorio.", RED),
        ("Zabbix 6.x", "Agentes, API, triggers, dashboards, histórico y correo integrados.", "Más configuración inicial, pero reproducible con scripts.", TEAL),
    ]
    for i, (name, strength, limit, accent) in enumerate(alternatives):
        y = 1.72 + i * 1.05
        rect(s, 0.82, y, 2.35, 0.72, WHITE, LINE, True)
        text(s, name, 1.02, y + 0.16, 1.92, 0.23, 12, accent, True)
        rect(s, 3.38, y, 3.9, 0.72, (245, 248, 250), LINE, True)
        text(s, strength, 3.58, y + 0.14, 3.45, 0.28, 10.5, INK)
        rect(s, 7.55, y, 4.55, 0.72, (250, 247, 244), LINE, True)
        text(s, limit, 7.75, y + 0.14, 4.05, 0.28, 10.5, INK)
    label_box(s, "Decisión", "Zabbix permite cumplir el enunciado sin SaaS pago: monitoreo de hosts, servicios, triggers, dashboards, histórico, alertas y API en Docker Compose.", 0.85, 6.03, 11.25, 0.78, TEAL)
    footer(s, 15)

    # 16 discussion
    s = prs.slides.add_slide(blank)
    bg(s, (242, 247, 249))
    title(s, "discusión", "La prueba importante no es instalar Zabbix; es cerrar el ciclo operativo.", 16)
    label_box(s, "Qué probamos", "Carga real con Artillery, caída controlada, detección en Zabbix, correo MailHog, recuperación y métricas históricas.", 0.85, 1.9, 3.55, 1.42, TEAL)
    label_box(s, "Qué aprendimos", "Monitorear no es solo saber si un contenedor vive: también hay que medir servicio, backend, DB, SLO y rutas bajo presión.", 4.85, 1.9, 3.55, 1.42, CYAN)
    label_box(s, "Limitación honesta", "El SLO del portal es de laboratorio y vive en memoria del proceso; Zabbix conserva el histórico real de eventos e items.", 8.85, 1.9, 3.55, 1.42, AMBER)
    text(s, "Conclusión", 0.9, 4.18, 2.0, 0.25, 11, MUTED, True)
    text(s, "El proyecto cumple los requisitos base y los amplía con una aplicación real observable: frontend, backend, MariaDB, /metrics, SLO, analíticas, Artillery y matriz 13/13.", 0.9, 4.62, 10.8, 0.82, 24, INK, True)
    footer(s, 16)

    # 17 references
    s = prs.slides.add_slide(blank)
    bg(s, INK)
    title(s, "referencias", "Las decisiones se apoyan en documentación oficial y herramientas reproducibles.", 17, True)
    refs = [
        ("Zabbix Documentation 6.0", "zabbix.com/documentation/6.0", TEAL),
        ("Docker Compose", "docs.docker.com/compose", CYAN),
        ("MailHog", "github.com/mailhog/MailHog", AMBER),
        ("Artillery Docs", "artillery.io/docs", LIME),
        ("PostgreSQL / MariaDB", "postgresql.org/docs + mariadb.org/documentation", CYAN),
        ("Caddy Automatic HTTPS", "caddyserver.com/docs/automatic-https", TEAL),
    ]
    for i, (name, url, accent) in enumerate(refs):
        x = 0.9 + (i % 2) * 5.75
        y = 1.85 + (i // 2) * 1.18
        rect(s, x, y, 5.05, 0.76, (16, 40, 51), (48, 73, 84), True)
        text(s, name, x + 0.2, y + 0.14, 3.1, 0.22, 11.5, accent, True)
        text(s, url, x + 0.2, y + 0.43, 4.55, 0.18, 8.5, (202, 213, 219))
    text(s, "Entregables: informe IEEE, presentación, repositorio GitHub, README, scripts de aprovisionamiento, Docker Compose, evidencias y demo pública HTTPS.", 0.95, 5.82, 10.9, 0.5, 17, WHITE, True)
    footer(s, 17, True)

    # Reorder the narrative so the deck follows the report structure:
    # problema/contexto -> alternativas -> objetivo/diseño -> implementación
    # -> pruebas -> discusión/conclusión -> entregables/demo -> referencias.
    logical_order = [0, 1, 14, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 15, 12, 13, 16]
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        prs.slides._sldIdLst.remove(slide_id)
    for idx in logical_order:
        prs.slides._sldIdLst.append(slide_ids[idx])

    dark_slides = {1, 8, 11, 16, 17}
    for new_num, slide in enumerate(prs.slides, start=1):
        dark = new_num in dark_slides
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            value = shape.text.strip()
            if value.endswith(f"/{total_slides}") and len(value) <= 6:
                tf = shape.text_frame
                tf.clear()
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.RIGHT
                r = p.add_run()
                r.text = f"{new_num:02d}/{total_slides}"
                r.font.name = "Aptos"
                r.font.size = PptPt(7.5)
                r.font.color.rgb = rgb((142, 159, 170) if dark else MUTED)
            elif value in {f"{i:02d}" for i in range(1, total_slides + 1)} and shape.top < PptInches(1) and shape.left < PptInches(1):
                tf = shape.text_frame
                tf.clear()
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = f"{new_num:02d}"
                r.font.name = "Aptos"
                r.font.size = PptPt(8)
                r.font.bold = True
                r.font.color.rgb = rgb(TEAL if dark else INK)

    out = OUT / "Presentacion_Proyecto7_Zabbix.pptx"
    prs.save(out)
    return out


def create_pptx_pdf(pptx_path):
    pdf_path = OUT / "Presentacion_Proyecto7_Zabbix.pdf"
    soffice = shutil.which("soffice") or shutil.which("soffice.com")
    if not soffice:
        return pdf_path if pdf_path.exists() else None
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(OUT), str(pptx_path)],
        check=False,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    return pdf_path if pdf_path.exists() else None


def create_checklist():
    path = OUT / "CHECKLIST_ENTREGA.md"
    path.write_text(
        """# Checklist de entrega - Proyecto 7

- [x] docker-compose.yml con Zabbix Server, Zabbix Web, base de datos, MailHog y 4 servicios monitoreados.
- [x] Imagen Zabbix personalizada con Dockerfile.
- [x] Archivos de configuración Zabbix montados como volumen.
- [x] Hosts monitoreados: web, DB, DNS y FTP.
- [x] Triggers de disponibilidad configurados.
- [x] Alertas por correo validadas en MailHog.
- [x] README con instrucciones de despliegue.
- [x] Informe IEEE en DOCX y PDF.
- [x] Presentacion PPTX.
- [x] Evidencias PNG de dashboard, hosts, latest data, falla y MailHog.
- [x] Web pública con backend, MariaDB, gráficas, SLO, exporter `/metrics` y pruebas Artillery.
- [x] Matriz de cumplimiento verificable en `/api/compliance`.
- [x] Script de auditoría reproducible `scripts/audit-project.sh`.
- [x] Guía de entregables y evaluación `docs/ENTREGABLES_EVALUACION.md`.
- [x] Matriz de rúbrica `docs/MATRIZ_RUBRICA.md`.

Accesos locales:

- Zabbix: http://localhost:8088
- Usuario: Admin
- Clave: zabbix
- MailHog: http://localhost:8025

Accesos públicos:

- Portal: https://web-zabbix.negociocontigo.com
- Zabbix: https://zabbix.negociocontigo.com
- MailHog: https://mailhog-zabbix.negociocontigo.com/login
""",
        encoding="utf-8",
    )
    return path


def zip_delivery(paths):
    zip_path = OUT / "Proyecto7_Zabbix_entrega_final.zip"
    include = [
        ROOT / ".env.example",
        ROOT / "docker-compose.yml",
        ROOT / "docker-compose.vps.yml",
        ROOT / "README.md",
        ROOT / "docker" / "zabbix-server" / "Dockerfile",
        ROOT / "docker" / "zabbix-server" / "zabbix_server.conf.d" / "proyecto7.conf",
        ROOT / "zabbix-config" / "agent" / "proyecto7-agent.conf",
        ROOT / "docs" / "PRUEBAS.md",
        ROOT / "docs" / "SUSTENTACION.md",
        ROOT / "docs" / "PRESENTACION.md",
        ROOT / "docs" / "DEMO_AVANZADA.md",
        ROOT / "docs" / "ENTREGABLES_EVALUACION.md",
        ROOT / "docs" / "MATRIZ_RUBRICA.md",
        ROOT / "docs" / "GUION_SUSTENTACION_20_MIN.md",
        ROOT / "entrega-final" / "GUION_PARA_ENVIAR_GRUPO.md",
        ROOT / "scripts" / "provision.ps1",
        ROOT / "scripts" / "provision_zabbix.py",
        ROOT / "scripts" / "test-failure.ps1",
        ROOT / "scripts" / "verify.ps1",
        ROOT / "scripts" / "audit-project.sh",
        ROOT / "scripts" / "evidence-pack.sh",
        ROOT / "scripts" / "live-artillery-demo.sh",
        ROOT / "scripts" / "demo-full.sh",
        ROOT / "tests" / "artillery-smoke.yml",
        ROOT / "tests" / "artillery-live-demo.yml",
        ROOT / "tests" / "artillery-web-service.yml",
        ROOT / "tests" / "artillery-stress-demo.yml",
        ROOT / "services" / "web" / "default.conf",
        ROOT / "services" / "web" / "Dockerfile",
        ROOT / "services" / "web" / "server.js",
        ROOT / "services" / "web" / "package.json",
        ROOT / "services" / "web" / "package-lock.json",
        ROOT / "services" / "web" / "html" / "index.html",
        ROOT / "services" / "dns" / "Corefile",
        *paths,
        *sorted(EVID.glob("*.png")),
    ]
    with ZipFile(zip_path, "w", ZIP_DEFLATED) as zf:
        for path in include:
            if path.exists():
                zf.write(path, path.relative_to(ROOT))
    return zip_path


def main():
    ensure_dirs()
    docx = create_docx()
    pdf = create_pdf()
    pptx = create_pptx()
    pptx_pdf = create_pptx_pdf(pptx)
    checklist = create_checklist()
    delivery_paths = [docx, pdf, pptx, checklist]
    if pptx_pdf:
        delivery_paths.append(pptx_pdf)
    zip_path = zip_delivery(delivery_paths)
    for path in [*delivery_paths, zip_path]:
        print(path)


if __name__ == "__main__":
    main()
